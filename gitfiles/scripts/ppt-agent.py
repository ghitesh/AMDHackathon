import json
from dataclasses import dataclass
from typing import List

from openai import OpenAI
from pptx import Presentation


# ============================================================================
# Data Models
# ============================================================================

@dataclass
class Slide:
    title: str
    bullet_points: List[str]
    speaker_notes: str


@dataclass
class PresentationContent:
    slides: List[Slide]


# ============================================================================
# JSON Schema for Guided Decoding
# ============================================================================

PRESENTATION_JSON_SCHEMA = {
    "type": "object",
    "properties": {
        "slides": {
            "type": "array",
            "items": {
                "type": "object",
                "properties": {
                    "title": {
                        "type": "string"
                    },
                    "bullet_points": {
                        "type": "array",
                        "items": {
                            "type": "string"
                        }
                    },
                    "speaker_notes": {
                        "type": "string"
                    }
                },
                "required": [
                    "title",
                    "bullet_points",
                    "speaker_notes"
                ]
            }
        }
    },
    "required": [
        "slides"
    ]
}


# ============================================================================
# OpenAI/vLLM Client
# ============================================================================

client = OpenAI(
    base_url="http://localhost:8000/v1",
    api_key="not-needed"
)

try:
    LOCAL_MODEL = client.models.list().data[0].id
except Exception:
    print("Could not connect to vLLM.")
    LOCAL_MODEL = "default-model"


# ============================================================================
# Parser
# ============================================================================

def parse_presentation(data: dict) -> PresentationContent:
    """Convert JSON dictionary into dataclass objects."""

    slides = []

    for slide in data.get("slides", []):

        slides.append(
            Slide(
                title=slide["title"],
                bullet_points=slide.get("bullet_points", []),
                speaker_notes=slide.get("speaker_notes", "")
            )
        )

    return PresentationContent(slides=slides)


# ============================================================================
# Content Agent
# ============================================================================

def content_agent(
        user_input: str,
        num_slides: int
) -> PresentationContent:

    print(
        f"Generating {num_slides} slides "
        f"using model '{LOCAL_MODEL}'..."
    )

    prompt = f"""
You are an expert presentation designer, technical writer, and educator.

Your task is to create a professional, presentation-ready slide deck.

Topic:
{user_input}

Number of Slides:
{num_slides}

Requirements:

- Create exactly {num_slides} slides.
- Organize the presentation logically from introduction to conclusion.
- Give every slide a concise, meaningful title.
- Include 3 to 6 concise bullet points per slide.
- Keep bullet points presentation-friendly.
- Avoid redundancy across slides.
- Expand short inputs using common domain knowledge.
- Summarize long inputs while preserving important information.
- Generate 2 to 5 sentences of speaker notes that expand on the slide.
- Make the presentation suitable for a professional audience.
"""

    response = client.chat.completions.create(

        model=LOCAL_MODEL,

        messages=[

            {
                "role": "system",
                "content": """
You are an expert presentation generation assistant.

Always respond with exactly one JSON object that conforms to the supplied JSON schema.

Do NOT generate:

- Markdown
- Code fences
- Explanations
- Introductory text
- Closing remarks
- Comments

Generate only professional presentation content.
"""
            },

            {
                "role": "user",
                "content": prompt
            }

        ],

        temperature=0.3,

        extra_body={
            "guided_json": PRESENTATION_JSON_SCHEMA
        }

    )

    raw = response.choices[0].message.content.strip()

    try:
        data = json.loads(raw)
    except json.JSONDecodeError:
        print(raw)
        raise

    return parse_presentation(data)


# ============================================================================
# PowerPoint Export Agent
# ============================================================================

def ppt_export_agent(
        presentation: PresentationContent,
        output_filename: str = "output.pptx"
):

    if not presentation.slides:
        raise ValueError("No slides generated.")

    prs = Presentation()

    # ----------------------------------------------------------------------
    # Title Slide
    # ----------------------------------------------------------------------

    title_layout = prs.slide_layouts[0]

    slide = prs.slides.add_slide(title_layout)

    slide.shapes.title.text = presentation.slides[0].title

    subtitle = slide.placeholders[1]

    if presentation.slides[0].bullet_points:
        subtitle.text = presentation.slides[0].bullet_points[0]

    # ----------------------------------------------------------------------
    # Content Slides
    # ----------------------------------------------------------------------

    for slide_data in presentation.slides[1:]:

        layout = prs.slide_layouts[1]

        slide = prs.slides.add_slide(layout)

        slide.shapes.title.text = slide_data.title

        body = slide.placeholders[1]

        tf = body.text_frame

        for i, bullet in enumerate(slide_data.bullet_points):

            if i == 0:
                tf.text = bullet

            else:
                p = tf.add_paragraph()
                p.text = bullet
                p.level = 0

        notes = slide.notes_slide.notes_text_frame

        notes.text = slide_data.speaker_notes

    prs.save(output_filename)

    print(f"Presentation saved as {output_filename}")


# ============================================================================
# Main
# ============================================================================

def main():

    topic = input(
        "Enter presentation topic: "
    )

    slide_count = int(
        input(
            "Number of slides: "
        )
    )

    presentation = content_agent(
        topic,
        slide_count
    )

    ppt_export_agent(
        presentation,
        "Local_Agent_Presentation.pptx"
    )


if __name__ == "__main__":
    main()
