from __future__ import annotations

from pathlib import Path
from typing import Optional

from pptx import Presentation
from pptx.enum.shapes import (
    MSO_AUTO_SHAPE_TYPE,
    MSO_CONNECTOR,
)
from pptx.enum.text import (
    MSO_ANCHOR,
    PP_ALIGN,
)
from pptx.dml.color import RGBColor
from pptx.util import Pt, Emu


class PowerPointRenderer:
    """
    Render DiagramModel into PowerPoint.

    Input coordinates must already be normalized
    into PowerPoint EMU space.
    """

    def __init__(
        self,
        prs: Presentation,
        icon_registry=None,
    ):
        self.prs = prs
        self.icon_registry = icon_registry
        self.node_shapes = {}

    # =====================================================
    # Public API
    # =====================================================

    def render(self, diagram):

        slide = self.prs.slides.add_slide(
            self.prs.slide_layouts[6]
        )

        self._render_containers(
            slide,
            diagram,
        )

        self._render_nodes(
            slide,
            diagram,
        )

        self._render_edges(
            slide,
            diagram,
        )

        return slide

    # =====================================================
    # Containers
    # =====================================================

    def _render_containers(
        self,
        slide,
        diagram,
    ):

        for container in diagram.containers:

            bbox = container.bbox

            shape = slide.shapes.add_shape(
                MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
                int(bbox.x),
                int(bbox.y),
                int(bbox.width),
                int(bbox.height),
            )

            shape.fill.background()

            shape.line.color.rgb = RGBColor(
                180,
                180,
                180,
            )

            shape.line.width = 12700

            # PowerPoint's default rounded-rectangle corner radius
            # is a fixed fraction (~1/6) of the shorter side. For
            # large containers that radius becomes huge in absolute
            # terms and can visually clip a node sitting near a
            # corner. AWS reference diagrams also use a much
            # subtler corner radius than PowerPoint's default.
            try:
                shape.adjustments[0] = 0.04
            except Exception:
                pass

            self._render_container_label(
                shape,
                container,
            )

    def _render_container_label(
        self,
        shape,
        container,
    ):

        # AWS-style group labels live in the top-left corner
        # of the box. python-pptx defaults to vertically and
        # horizontally centered text, which is what caused
        # nested container labels (region / vpc / az / subnet)
        # to all pile up near the same point and overlap with
        # the node icons sitting on top of them.

        text_frame = shape.text_frame

        text_frame.text = container.label

        text_frame.word_wrap = True

        text_frame.vertical_anchor = MSO_ANCHOR.TOP

        text_frame.margin_left = Emu(91440)
        text_frame.margin_top = Emu(45720)
        text_frame.margin_right = Emu(0)
        text_frame.margin_bottom = Emu(0)

        paragraph = text_frame.paragraphs[0]

        paragraph.alignment = PP_ALIGN.LEFT

        for run in paragraph.runs:

            run.font.size = Pt(10)

            run.font.color.rgb = RGBColor(
                107,
                114,
                128,
            )

    # =====================================================
    # Nodes
    # =====================================================

    def _render_nodes(
        self,
        slide,
        diagram,
    ):
    
        for node in diagram.nodes:
    
            # Skip Graphviz helper nodes
    
            if (
                node.id.startswith("layer_")
                or node.id.startswith("align_")
                or ";" in node.id
            ):
                continue
    
            print(
                "ppt_generator:105 node.service=",
                getattr(node, "service", None)
            )
    
            icon_path = self._resolve_icon(
                getattr(node, "service", None)
            )
    
            left = int(
                node.x - (node.width / 2)
            )
    
            top = int(
                node.y - (node.height / 2)
            )
    
            width = int(node.width)
            height = int(node.height)
    
            if icon_path:
                print("ppt_generator:105 icon.path=", icon_path)
                shape = slide.shapes.add_picture(
                    str(icon_path),
                    left,
                    top,
                    width=width,
                    height=height,
                )
    
            else:
    
                shape = slide.shapes.add_shape(
                    MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
                    left,
                    top,
                    width,
                    height,
                )
    
            self.node_shapes[node.id] = shape
    
            self._render_node_label(
                slide,
                node,
            )

    def normalize_service_name(self, name: Optional[str]) -> str:
        name = name.lower().strip()

        if name in {"alb", "application load balancer", "application-load-balancer"}:
            return "ALB"

        if name in {"nlb", "network load balancer"}:
            return "NLB"

        return name.upper()

    def _resolve_icon(
        self,
        service: Optional[str],
    ):
    
        if (
            not service
            or not self.icon_registry
        ):
            print(
                "ppt_generator.py:229 missing service or registry:",
                service
            )
            return None

        service = self.normalize_service_name(service)

        try:
    
            path = self.icon_registry.get_icon(
                service
            )
    
            print(
                f"ICON LOOKUP service={service} path={path}"
            )
    
            if (
                path
                and Path(path).exists()
            ):
                print(
                    f"ICON FOUND service={service}"
                )
                return path
    
            print(
                f"ICON MISSING service={service} path={path}"
            )
    
        except Exception as ex:
    
            print(
                f"ICON LOOKUP ERROR service={service}: {ex}"
            )
    
        return None

    def _render_node_label(
        self,
        slide,
        node,
    ):
    
        label_height = 250000
    
        left = int(
            node.x - (node.width / 2)
        )
    
        top = int(
            node.y + (node.height / 2)
        )
    
        width = int(
            max(node.width, 800000)
        )
    
        textbox = slide.shapes.add_textbox(
            left,
            top,
            width,
            label_height,
        )
    
        textbox.text = node.label
    
    # =====================================================
    # Edges
    # =====================================================

    def _render_edges(
        self,
        slide,
        diagram,
    ):

        for edge in diagram.edges:

            points = getattr(
                edge,
                "points",
                None,
            )

            if not points:
                continue

            if len(points) < 2:
                continue

            for p1, p2 in zip(
                points,
                points[1:],
            ):

                try:

                    connector = (
                        slide.shapes.add_connector(
                            MSO_CONNECTOR.STRAIGHT,
                            int(p1.x),
                            int(p1.y),
                            int(p2.x),
                            int(p2.y),
                        )
                    )

                    connector.line.width = 12700

                except Exception as ex:

                    print(
                        f"Edge render failed "
                        f"{edge.source}->{edge.target}: "
                        f"{ex}"
                    )

    # =====================================================
    # Save helper
    # =====================================================

    def save(
        self,
        output_file,
    ):

        self.prs.save(
            str(output_file)
        )


# =========================================================
# Backwards Compatibility
# =========================================================

def build_ppt(
    prs: Presentation,
    icon_registry=None,
):

    return PowerPointRenderer(
        prs,
        icon_registry,
    )
